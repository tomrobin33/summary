from pptx import Presentation

def parse_ppt(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    texts.append(shape.text)
        slides.append({"slide": i+1, "texts": texts})
    return slides 