import os
import tempfile
import requests
from fastapi import FastAPI, Query
from pydantic import BaseModel
from typing import Optional

# 文件解析库
import pandas as pd
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation

app = FastAPI()

class ParseRequest(BaseModel):
    url: str
    file_type: Optional[str] = None  # excel, word, ppt

@app.post("/parse_file")
def parse_file(req: ParseRequest):
    # 下载文件
    response = requests.get(req.url)
    if response.status_code != 200:
        return {"error": "下载失败"}
    # 临时保存
    suffix = ''
    if req.file_type == 'excel':
        suffix = '.xlsx'
    elif req.file_type == 'word':
        suffix = '.docx'
    elif req.file_type == 'ppt':
        suffix = '.pptx'
    else:
        # 自动判断
        if req.url.endswith('.xlsx'):
            req.file_type = 'excel'
            suffix = '.xlsx'
        elif req.url.endswith('.docx'):
            req.file_type = 'word'
            suffix = '.docx'
        elif req.url.endswith('.pptx'):
            req.file_type = 'ppt'
            suffix = '.pptx'
        else:
            return {"error": "无法识别文件类型"}
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(response.content)
        tmp_path = tmp.name
    # 解析内容
    try:
        if req.file_type == 'excel':
            content = parse_excel(tmp_path)
        elif req.file_type == 'word':
            content = parse_word(tmp_path)
        elif req.file_type == 'ppt':
            content = parse_ppt(tmp_path)
        else:
            content = []
    finally:
        os.remove(tmp_path)
    return {"file_type": req.file_type, "content": content}

def parse_excel(path):
    wb = load_workbook(path, data_only=True)
    result = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append(list(row))
        result.append({"sheet": sheet, "data": data})
    return result

def parse_word(path):
    doc = Document(path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append({"type": "paragraph", "text": para.text})
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            table_data.append([cell.text for cell in row.cells])
        content.append({"type": "table", "data": table_data})
    return content

def parse_ppt(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    texts.append(shape.text)
        slides.append({"slide": i+1, "texts": texts})
    return slides 